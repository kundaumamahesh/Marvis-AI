import json
import re

from pptx import Presentation

from office.powerpoint.slide_factory import (
    SlideFactory
)


class PPTGenerator:

    async def generate(
        self,
        structured_json,
        output
    ):

        prs = Presentation()

        try:

            # --------------------------------
            # CLEAN RAW MODEL RESPONSE
            # --------------------------------
            cleaned = structured_json.strip()

            # remove markdown wrappers
            cleaned = re.sub(
                r"```json|```",
                "",
                cleaned,
                flags=re.IGNORECASE
            ).strip()

            # extract JSON safely
            match = re.search(
                r"\{.*\}",
                cleaned,
                re.DOTALL
            )

            if not match:
                raise Exception(
                    "No valid JSON found from AI."
                )

            data = json.loads(
                match.group()
            )

        except Exception as e:

            print(
                "[PPT JSON ERROR]",
                str(e)
            )

            # fallback PPT
            slide = prs.slides.add_slide(
                prs.slide_layouts[1]
            )

            slide.shapes.title.text = (
                "MARVIS PowerPoint Error"
            )

            slide.placeholders[1].text = (
                "The AI failed to generate "
                "valid slide JSON."
            )

            prs.save(output)
            return output

        # --------------------------------
        # SLIDES SAFETY
        # --------------------------------
        slides = data.get(
            "slides",
            []
        )

        if not slides:

            slide = prs.slides.add_slide(
                prs.slide_layouts[1]
            )

            slide.shapes.title.text = (
                "No Slides Generated"
            )

            slide.placeholders[1].text = (
                "MARVIS did not receive "
                "valid slide content."
            )

            prs.save(output)
            return output

        # --------------------------------
        # BUILD SLIDES
        # --------------------------------
        for slide_data in slides:

            slide_type = (
                slide_data
                .get("type", "")
                .lower()
                .strip()
            )

            try:

                # TITLE
                if slide_type == "title":

                    SlideFactory.title_slide(
                        prs,
                        slide_data
                    )

                # BULLETS
                elif slide_type in [
                    "bullet",
                    "content",
                    "text"
                ]:

                    SlideFactory.bullet_slide(
                        prs,
                        slide_data
                    )

                # IMAGE SLIDE
                elif slide_type == "image":

                    SlideFactory.image_slide(
                        prs,
                        slide_data
                    )

                # FALLBACK
                else:

                    SlideFactory.bullet_slide(
                        prs,
                        {
                            "title":
                            slide_data.get(
                                "title",
                                "Slide"
                            ),

                            "content":
                            slide_data.get(
                                "content",
                                [
                                    str(slide_data)
                                ]
                            )
                        }
                    )

            except Exception as e:

                print(
                    "[PPT SLIDE ERROR]",
                    str(e)
                )

        # --------------------------------
        # SAVE FILE
        # --------------------------------
        prs.save(output)

        return output