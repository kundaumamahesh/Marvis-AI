import os
from pptx.util import Pt, Inches


class SlideFactory:

    @staticmethod
    def title_slide(
        prs,
        slide_data
    ):

        slide_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(
            slide_layout
        )

        title = slide.shapes.title

        subtitle = slide.placeholders[1]

        title.text = slide_data["title"]

        subtitle.text = slide_data["subtitle"]

        title.text_frame.paragraphs[0].font.size = Pt(30)

    @staticmethod
    def bullet_slide(
        prs,
        slide_data
    ):

        slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(
            slide_layout
        )

        title = slide.shapes.title

        title.text = slide_data["title"]

        body = slide.placeholders[1]

        tf = body.text_frame

        bullets = slide_data.get("bullets") or slide_data.get("bullets") or []
        if not bullets and "content" in slide_data:
            bullets = slide_data["content"]
            if isinstance(bullets, str):
                bullets = [bullets]

        if bullets:
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
        else:
            tf.text = ""

    @staticmethod
    def image_slide(
        prs,
        slide_data
    ):
        # Use Title Only layout (index 5)
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title = slide.shapes.title
        title.text = slide_data.get("title", "Visual Highlight")

        # Extract image path
        image_path = slide_data.get("image_path") or slide_data.get("image") or slide_data.get("url")
        
        if image_path and os.path.exists(image_path):
            left = Inches(1.5)
            top = Inches(2.0)
            width = Inches(7.0)
            slide.shapes.add_picture(image_path, left, top, width=width)
        else:
            # Fallback if image not found on disk
            left = Inches(1.5)
            top = Inches(3.0)
            width = Inches(7.0)
            height = Inches(2.0)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[Image Visual Placeholder: {image_path or 'No path specified'}]"
            p.font.size = Pt(18)
            p.font.italic = True

